import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from tabulate import tabulate
from tqdm import tqdm
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Pt
import re
import tkinter as tk
from tkinter import messagebox
from tkinter import ttk
import pygame

def generate_report(image_path,ppt_path):

    def delete_empty_shapes(presentation_path):
        try:
            prs = Presentation(presentation_path)
            for slide in prs.slides:
                slide_shapes = slide.shapes
                for shape in slide_shapes:
                    if shape.has_text_frame and not shape.text_frame.text:
                        slide_shapes._spTree.remove(shape._element)
            prs.save(presentation_path)
        except Exception as e:
            print(f"An error occurred during shape deletion: {e}")


    def extract_number(filename):
        match = re.search(r'(\d+)Nm', filename)
        if match:
            return int(match.group(1))
        else:
            return 0
    def order_number(filename):
        match = re.search(r'(\d+)order', filename)
        if match:
            return int(match.group(1))
        else:
            return 0
    def ppt_export(image_dir,ppt_path):
        # 读取文本文件内容
        with open('data\\Project\\ProcessedData\\result.txt', 'r') as file:
            text_lines = file.readlines()
        ppt = Presentation(ppt_path)
        subdirs = [d for d in os.listdir(image_dir) if os.path.isdir(os.path.join(image_dir, d))]
        image_files = []
        
        if not subdirs:
            
            image_files = [file for file in os.listdir(image_dir) if file.endswith(('.jpg', '.png', '.jpeg', '.gif'))]
            sorted_image_files = sorted(image_files, key=extract_number)
            running_up_files = [f for f in sorted_image_files if 'Running Up' in f]
            running_down_files = [f for f in sorted_image_files if 'Running Down' in f]
            sorted_image_files = running_up_files + running_down_files

            for i, file in enumerate(sorted_image_files):
                slide = ppt.slides.add_slide(ppt.slide_layouts[8])
                slide.shapes.title.text = os.path.splitext(file)[0]
                image_file = os.path.join(image_dir, file)
                left = Inches(0.4)
                top = Inches(1)
                width = Inches(6)
                height = Inches(4)
                slide.shapes.add_picture(image_file, left, top, width, height)

            ppt.save(ppt_path)

        else:
            for subdir in subdirs:
                if subdir == 'avg':
                    subdir_path = []
                    subdir_path = os.path.join(image_dir, subdir)
                    print(subdir_path)
                    image_files =[]
                    image_files = [f for f in os.listdir(subdir_path) if f.endswith('.jpg') or f.endswith('.png')]

                    print(image_files)
                    

                    sorted_image_files = sorted(image_files, key=lambda x: (extract_number(x), order_number(x)))

                    text_lines.sort(key=lambda x: (extract_number(x)))

                    print(subdir)
                
                    for i, file in enumerate(sorted_image_files):
                        slide = ppt.slides.add_slide(ppt.slide_layouts[8])
                        slide.shapes.title.text = os.path.splitext(file)[0]
                        image_file = os.path.join(subdir_path, file)
                        left = Inches(0.4)
                        top = Inches(1)
                        width = Inches(5.6)
                        height = Inches(4)
                        slide.shapes.add_picture(image_file, left, top, width, height)

                
                        text_start = i * 8+1
                        text_end = min((i + 1) * 8+1, len(text_lines)-1)
                        text_content = ''.join(text_lines[text_start:text_end])


                        left = Inches(6)
                        top = Inches(0.7)
                        width = Inches(3.6)
                        height = Inches(4)
                        textbox = slide.shapes.add_textbox(left, top, width, height)
                        text_frame = textbox.text_frame

                        p = text_frame.add_paragraph()
                        p.text = text_content

                        # 设置加粗文本样式
                        p.font.bold = False
                        p.font.size = Pt(10)
                        p.alignment = PP_ALIGN.LEFT


                        # 设置文本自动换行
                        text_frame.word_wrap = True
                        textbox.text_frame.auto_size = True
                        textbox.text_frame.wrap_text = True
                else:
                    subdir_path = []
                    subdir_path = os.path.join(image_dir, subdir)
                    print(subdir_path)
                    image_files =[]
                    image_files = [f for f in os.listdir(subdir_path) if f.endswith('.jpg') or f.endswith('.png')]

                    print(image_files)
                    

                    sorted_image_files = sorted(image_files, key=lambda x: (extract_number(x), order_number(x)))

                    text_lines.sort(key=lambda x: (extract_number(x)))

                    print(subdir)
                
                    for i, file in enumerate(sorted_image_files):
                        slide = ppt.slides.add_slide(ppt.slide_layouts[8])
                        slide.shapes.title.text = os.path.splitext(file)[0]
                        image_file = os.path.join(subdir_path, file)
                        left = Inches(0.4)
                        top = Inches(1)
                        width = Inches(5.6)
                        height = Inches(4)
                        slide.shapes.add_picture(image_file, left, top, width, height)


        ppt.save(ppt_path)

                


    ppt_export(image_path,ppt_path)

    delete_empty_shapes(ppt_path)

    root = tk.Tk()
    root.withdraw()  # 隐藏主窗口

    messagebox.showinfo("提示", "Congratulation!报告完成")
    def play_sound():
                pygame.mixer.init()
                pygame.mixer.music.load("gui//sound.wav")
                pygame.mixer.music.play()

    play_sound()

    root.destroy()  # 销毁主窗口



